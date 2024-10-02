import docx
from docx import Document
from pptx import Presentation
import PyPDF2

import os
class pdf:
    def __init__(self,path) -> None:
        self.path = path
    def pdf_to_text(self):
        # Open the PDF file in read-binary mode
        os.chdir(f'{self.path}')
        files = os.listdir()
        pdf = list()
        for file in files: 
            if '.pdf' in file:
                pdf.append(file)
        print(pdf)
        all_pdf_texts = list()
        for i in pdf:

            with open(self.path + "\\" + f"{i}", 'rb') as pdf_file:
                # Create a PdfReader object instead of PdfFileReader
                pdf_reader = PyPDF2.PdfReader(pdf_file)

                text = ''

                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text()

            all_pdf_texts.append(text)
        return all_pdf_texts
    

class pptx_files:
    def __init__(self, path): 
        self.path = path

    def pptx_to_text(self):
        os.chdir(f'{self.path}')
        files = os.listdir()
        pptx_files = list()
        for ppt in files:
            if '.pptx' in ppt:
                pptx_files.append(ppt)
        print(pptx_files)
        ppt_text = list()
        for i in pptx_files:
            prs = Presentation(self.path + "\\" + f"{i}")
            print(i)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            text_runs = ",".join(text_runs)
            ppt_text.append(text_runs)
        return ppt_text

# Converts docx to text by a single function call. 
class DocxtoText: 
    def __init__(self,path) -> None:
        self.path = path 
    def docx_to_text(self):
        os.chdir(self.path)
        files = os.listdir()
        docx_files = list()
        for i in files:
            if '.docx' in i:
                docx_files.append(i)
        print(docx_files)     
        texts = list()
        for i in docx_files:
            doc = Document(self.path + r"\\" + f"{i}")
            full_text = []
            for paragraph in doc.paragraphs:
                full_text.append(paragraph.text)
            texts.append('\n'.join(full_text))
        return texts
a = DoctoText(path = r"C:\Users\Ibrahim.intern\Desktop\Converter")
text = a.docx_to_text()
print(len(text))


