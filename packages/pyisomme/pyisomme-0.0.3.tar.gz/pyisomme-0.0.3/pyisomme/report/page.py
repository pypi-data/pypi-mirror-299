from __future__ import annotations

from pyisomme import Channel, Isomme
from pyisomme.limits import limit_list_sort
from pyisomme.plotting import Plot_Line, Plot, Plot_Table, Plot_Line_Table
from pyisomme.report.criterion import Criterion
from pyisomme.unit import Unit, g0

import matplotlib.pyplot as plt
from matplotlib.colors import to_rgb
import numpy as np
import io
import os
from pptx.util import Inches
from abc import abstractmethod
from datetime import datetime
from typing import Callable


class Page:
    name: str

    def __init__(self, report):
        self.report = report

    @abstractmethod
    def construct(self, presentation) -> None:
        pass

    def __repr__(self):
        return f"Page({self.name})"


class Page_Cover(Page):
    name = "Cover"
    title: str = None
    subtitle: str = None

    def __init__(self, report):
        super().__init__(report)
        self.title = report.title
        self.subtitle = f'{report.name}\n{" | ".join([isomme.test_number for isomme in report.isomme_list])}'

    def construct(self, presentation):
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title
        slide.placeholders[1].text = self.subtitle


class Page_Content(Page):
    title: str = None
    footer: str = f"{datetime.now().strftime('%d.%m.%Y')} | {os.getlogin()}"

    def construct(self, presentation) -> None:
        title_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.title

        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        txBox = slide.shapes.add_textbox(0, Inches(slide_height / Inches(1) - 0.3), slide_width, Inches(0.3))
        tf = txBox.text_frame
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = self.footer
        p.font.size = Inches(0.2)


class Page_Criterion_Table(Page_Content):
    criteria: dict[Isomme, list[Criterion]]
    row_label: Callable
    cell_text: Callable

    def __init__(self, report):
        super().__init__(report)
        self.criteria = {}

    def construct(self, presentation):
        super().construct(presentation)
        slide = presentation.slides[-1]

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        figsize_y = 8
        figsize_x = figsize_y * float(width) / float(height)

        cell_text = np.full((len(list(self.criteria.values())[0]), len(list(self.criteria.keys()))), np.nan).tolist()
        cell_colors = np.zeros_like(cell_text).tolist()
        for idx_isomme, isomme in enumerate(self.criteria.keys()):
            for idx_criterion, criterion in enumerate(self.criteria[isomme]):
                cell_text[idx_criterion][idx_isomme] = self.cell_text(criterion)
                cell_colors[idx_criterion][idx_isomme] = (*to_rgb(criterion.color), 0.5) if criterion.color is not None else (0,0,0,0)

        row_labels = [self.row_label(criterion) for criterion in self.criteria[list(self.criteria.keys())[0]]]

        col_labels = [isomme.test_number for isomme in self.criteria.keys()]
        col_colors = [mcolor for mcolor in list(Plot.colors)[:len(self.criteria.keys())]]

        fig = Plot_Table(cell_texts=[cell_text],
                         cell_colors=[cell_colors],
                         row_labels=[row_labels],
                         col_labels=[col_labels],
                         col_labels_colors=[col_colors],
                         col_labels_fontweight="bold",
                         nrows=1,
                         ncols=1,
                         figsize=(figsize_x, figsize_y)).fig

        image_steam = io.BytesIO()
        fig.savefig(image_steam, transparent=True, bbox_inches='tight')
        slide.shapes.add_picture(image_steam, left=left, top=top, height=height)


class Page_Criterion_Values_Table(Page_Criterion_Table):
    row_label = staticmethod(lambda criterion: f"{criterion.name} [{criterion.channel.unit if criterion.channel is not None else np.nan}]")
    cell_text = staticmethod(lambda criterion: f"{criterion.value:.4g}")


class Page_Criterion_Rating_Table(Page_Criterion_Table):
    row_label = staticmethod(lambda criterion: f"{criterion.name}")
    cell_text = staticmethod(lambda criterion: f"{criterion.rating:.1f}")


class Page_Criterion_Values_Chart(Page_Content):
    criteria: dict[Isomme, list[Criterion]]

    def __init__(self, report):
        super().__init__(report)
        self.criteria = {}

    def construct(self, presentation):
        super().construct(presentation)
        slide = presentation.slides[-1]

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        figsize_y = 8
        figsize_x = figsize_y * float(width) / float(height)

        fig, ax = plt.subplots(figsize=(figsize_x, figsize_y), layout="constrained")

        bar_width = 0.8 / len(self.criteria)
        x_labels = [c.name for c in list(self.criteria.values())[0]]
        x = np.arange(len(x_labels))
        x_offsets = np.linspace(-0.4 + bar_width / 2, 0.4 - bar_width / 2, len(self.criteria))

        unique_limits = np.zeros(len(x_labels), dtype=bool)
        line_values = np.array([[abs(c.value) for c in criteria] for isomme, criteria in self.criteria.items()])

        # Same limits?
        for idx, c1 in enumerate(list(self.criteria.values())[0]):
            for c2_list in list(self.criteria.values())[1:]:
                c2 = c2_list[idx]
                if None in (c1.channel, c2.channel):
                    x_limit1 = x_limit2 = 0
                else:
                    x_limit1 = c1.limits.get_limit_min_x(c1.channel)
                    x_limit2 = c2.limits.get_limit_min_x(c2.channel)
                if not np.all([abs((l1.func(x_limit1) - l2.func(x_limit2)) / l1.func(x_limit1)) < 1e-6 for l1, l2 in zip(limit_list_sort(c1.limits.limit_list), limit_list_sort(c2.limits.limit_list))]):
                    unique_limits[idx] = True
                    break

        # Calculate Column Factor
        col_factors = np.nanmax(1.1 * np.abs(line_values), axis=0)

        for idx_isomme, criteria in enumerate(self.criteria.values()):
            for idx_col, criterion in enumerate(criteria):
                if criterion.channel is None:
                    x_limit = 0
                else:
                    x_limit = criterion.limits.get_limit_min_x(criterion.channel)

                col_factor_limit = 1.1 * np.nanmax([abs(l.func(x_limit)) if not np.isinf(abs(l.func(x_limit))) else np.nan for l in criterion.limits.limit_list])
                if not np.isnan(col_factor_limit):
                    col_factors[idx_col] = np.nanmax([col_factor_limit, col_factors[idx_col]])

        # Plot Bars
        for idx_isomme, (isomme, criteria) in enumerate(self.criteria.items()):
            for idx_col, criterion in enumerate(criteria):
                limits = limit_list_sort(criterion.limits.limit_list, sym=True)

                limit_values = []
                for idx, limit in enumerate(limits):
                    if criterion.channel is None:
                        x_limit = 0
                    else:
                        x_limit = criterion.limits.get_limit_min_x(criterion.channel)
                    limit_values.append(abs(limit.func(x_limit)) if not np.isinf(abs(limit.func(x_limit))) else col_factors[idx_col])

                for idx, (limit, limit_value) in enumerate(zip(limits, limit_values)):
                    if idx == 0:
                        bar_bottom = 0
                        bar_height = limit_value / col_factors[idx_col]
                    elif idx < len(limits) - 1:
                        if (limit.lower and limit.func(0) >= 0) or (limit.upper and limit.func(0) < 0):
                            bar_bottom = limit_value / col_factors[idx_col]
                            bar_height = (limit_values[idx+1] - limit_value) / col_factors[idx_col]
                        elif (limit.upper and limit.func(0) >= 0) or (limit.lower and limit.func(0) < 0):
                            bar_bottom = limit_values[idx-1] / col_factors[idx_col]
                            bar_height = (limit_value - limit_values[idx-1]) / col_factors[idx_col]
                        else:
                            continue
                    else:  # idx == len(limits) - 1
                        bar_bottom = limit_value / col_factors[idx_col]
                        bar_height = 1 - limit_value / col_factors[idx_col]

                    ax.bar(x=x[idx_col] + x_offsets[idx_isomme],
                           bottom=bar_bottom,
                           height=bar_height,
                           color=limit.color,
                           width=bar_width,
                           alpha=0.5,
                           label=limit.name)

        # Plot Lines
        for idx, isomme in enumerate(self.report.isomme_list):
            ax.plot(x + unique_limits * x_offsets[idx],
                    line_values[idx, :] / col_factors,
                    marker="o",
                    label=isomme.test_number,
                    linewidth=3,
                    markersize=8)

        ax.set_xticks(x, x_labels, rotation=30, ha='right')
        ax.get_yaxis().set_visible(False)

        # Legend (Delete duplicates)
        handles, labels = ax.get_legend_handles_labels()
        by_label = dict(zip(labels, handles))
        ax.legend(by_label.values(), by_label.keys(), bbox_to_anchor=(1, 1), loc='upper left')

        image_steam = io.BytesIO()
        fig.savefig(image_steam, transparent=True, bbox_inches='tight')
        slide.shapes.add_picture(image_steam, left=left, top=top, height=height)


class Page_Plot_nxn(Page_Content):
    channels: dict[Isomme, list[list[Channel | str]]]
    nrows: int = 1
    ncols: int = 1
    sharex: bool = False
    sharey: bool = False
    xlim: tuple[float | int, float | int] = None
    ylim: tuple[float | int, float | int] = None

    def __init__(self, report):
        super().__init__(report)
        if self.title is None:
            self.title = self.name

    def construct(self, presentation):
        super().construct(presentation)
        slide = presentation.slides[-1]

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        figsize_y = 8
        figsize_x = figsize_y * float(width) / float(height)

        fig = Plot_Line(self.channels,
                        nrows=self.nrows,
                        ncols=self.ncols,
                        sharex=self.sharex,
                        sharey=self.sharey,
                        xlim=self.xlim,
                        ylim=self.ylim,
                        limits=self.report.limits,
                        figsize=(figsize_x, figsize_y)).fig

        image_steam = io.BytesIO()
        fig.savefig(image_steam, transparent=True, bbox_inches='tight')
        slide.shapes.add_picture(image_steam, left=left, top=top, height=height)


class Page_Line_Table(Page_Content):
    channels: dict[Isomme, list[list[Channel | str]]]
    cell_texts: list[np.ndarray | list[list, ...]]
    row_labels: list[np.ndarray | list]
    col_labels: list[np.ndarray | list]
    cell_colors: list[np.ndarray | list[list, ...]] = None,
    col_labels_colors: list[np.ndarray | list] = None,
    col_labels_fontweight: str = None,
    nrows: int = 1
    ncols: int = 1
    sharex: bool = False
    sharey: bool = False
    xlim: tuple[float | int, float | int] = None
    ylim: tuple[float | int, float | int] = None

    def __init__(self, report):
        super().__init__(report)

    def construct(self, presentation):
        super().construct(presentation)
        slide = presentation.slides[-1]

        top = slide.placeholders[1].top
        left = slide.placeholders[1].left
        height = slide.placeholders[1].height
        width = slide.placeholders[1].width

        sp = slide.placeholders[1].element
        sp.getparent().remove(sp)

        figsize_y = 8
        figsize_x = figsize_y * float(width) / float(height)

        fig = Plot_Line_Table(channels=self.channels,
                              cell_texts=self.cell_texts,
                              row_labels=self.row_labels,
                              col_labels=self.col_labels,
                              nrows=self.nrows,
                              ncols=self.ncols,
                              sharex=self.sharex,
                              sharey=self.sharey,
                              xlim=self.xlim,
                              ylim=self.ylim,
                              limits=self.report.limits,
                              figsize=(figsize_x, figsize_y)).fig

        image_steam = io.BytesIO()
        fig.savefig(image_steam, transparent=True, bbox_inches='tight')
        slide.shapes.add_picture(image_steam, left=left, top=top, height=height)


class Page_OLC(Page_Line_Table):
    name = "OLC"
    title = "Occupant Load Criterion (OLC)"
    nrows: int = 1
    ncols: int = 2

    def __init__(self, report):
        super().__init__(report)

        self.channels = {isomme: [[isomme.get_channel("10VEHCCG00??VEXA", "14BPIL??????VEXA", "10SEATLERE??VEXA"),
                                   isomme.get_channel("10VEH0OLC??VEXA", "14BPIL0OLC??VEXA", "10SEAT0OLC??VEXA")]] for isomme in self.report.isomme_list}
        self.cell_texts = [[[f'{isomme.get_channel("10VEH0OLC??VEXX", "14BPIL0OLC??VEXX", "10SEAT0OLC??VEXX").get_data(unit=Unit(g0))[0] if isomme.get_channel("14BPIL0OLC??VEXX", "10SEAT0OLC??VEXX") is not None else np.nan:.2f}'] for isomme in self.report.isomme_list]]
        self.col_labels = [["OLC [g]"]]
        self.row_labels = [[isomme.test_number for isomme in self.report.isomme_list]]
